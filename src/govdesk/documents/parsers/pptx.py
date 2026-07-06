# SPDX-FileCopyrightText: 2026 GovDesk-Mitwirkende
#
# SPDX-License-Identifier: EUPL-1.2

"""PPTX-Parser (python-pptx): jede Folie wird zur Sektion mit ihrem Text."""

import io

from pptx import Presentation

from govdesk.documents.parsers.base import Block, ParsedDocument


class PptxParser:
    def parse(self, data: bytes) -> ParsedDocument:
        prs = Presentation(io.BytesIO(data))
        blocks: list[Block] = []
        for index, slide in enumerate(prs.slides, start=1):
            title = None
            if slide.shapes.title is not None:
                title = " ".join((slide.shapes.title.text or "").split())
            blocks.append(Block(text=title or f"Folie {index}", heading_level=2, page_no=index))
            for shape in slide.shapes:
                if shape == slide.shapes.title or not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    text = " ".join("".join(run.text for run in paragraph.runs).split())
                    if text:
                        blocks.append(Block(text=text, page_no=index))
        return ParsedDocument(blocks=blocks)
